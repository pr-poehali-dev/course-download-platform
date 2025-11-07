import json
import os
import io
import zipfile
import base64
from typing import Dict, Any, List
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import psycopg2

def handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    '''
    Business: Генерация пакета для защиты работы (презентация + речь + шпаргалка)
    Args: event - dict с httpMethod, body (form data)
          context - объект с request_id
    Returns: ZIP архив с файлами для защиты
    '''
    method: str = event.get('httpMethod', 'POST')
    
    if method == 'OPTIONS':
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'POST, OPTIONS',
                'Access-Control-Allow-Headers': 'Content-Type, X-Auth-Token',
                'Access-Control-Max-Age': '86400'
            },
            'body': '',
            'isBase64Encoded': False
        }
    
    if method != 'POST':
        return {
            'statusCode': 405,
            'headers': {'Content-Type': 'application/json', 'Access-Control-Allow-Origin': '*'},
            'body': json.dumps({'error': 'Метод не поддерживается'}),
            'isBase64Encoded': False
        }
    
    try:
        body_data = json.loads(event.get('body', '{}'))
        
        title = body_data.get('title', '').strip()
        relevance = body_data.get('relevance', '').strip()
        goal = body_data.get('goal', '').strip()
        tasks = [t.strip() for t in body_data.get('tasks', []) if t.strip()]
        conclusions = [c.strip() for c in body_data.get('conclusions', []) if c.strip()]
        university = body_data.get('university', '').strip()
        faculty = body_data.get('faculty', '').strip()
        author = body_data.get('author', '').strip()
        work_id = body_data.get('workId')
        user_id = body_data.get('userId')
        
        if not all([title, relevance, goal, tasks, conclusions, university, author]):
            return {
                'statusCode': 400,
                'headers': {'Content-Type': 'application/json', 'Access-Control-Allow-Origin': '*'},
                'body': json.dumps({'error': 'Заполните все обязательные поля'}),
                'isBase64Encoded': False
            }
        
        pptx_bytes = generate_presentation({
            'title': title,
            'relevance': relevance,
            'goal': goal,
            'tasks': tasks,
            'conclusions': conclusions,
            'university': university,
            'faculty': faculty,
            'author': author
        })
        
        docx_bytes = generate_speech({
            'title': title,
            'relevance': relevance,
            'goal': goal,
            'tasks': tasks,
            'conclusions': conclusions
        })
        
        notes_bytes = generate_notes({
            'title': title,
            'relevance': relevance,
            'goal': goal,
            'tasks': tasks,
            'conclusions': conclusions
        })
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            zip_file.writestr('презентация.pptx', pptx_bytes.getvalue())
            zip_file.writestr('речь_для_защиты.docx', docx_bytes.getvalue())
            zip_file.writestr('шпаргалка.txt', notes_bytes)
        
        zip_buffer.seek(0)
        zip_base64 = base64.b64encode(zip_buffer.read()).decode('utf-8')
        
        if work_id and user_id:
            save_defense_kit_record(user_id, work_id)
        
        return {
            'statusCode': 200,
            'headers': {
                'Content-Type': 'application/zip',
                'Access-Control-Allow-Origin': '*',
                'Content-Disposition': f'attachment; filename="defense_kit_{datetime.now().strftime("%Y%m%d_%H%M%S")}.zip"'
            },
            'body': zip_base64,
            'isBase64Encoded': True
        }
        
    except Exception as e:
        print(f"Defense kit generation error: {repr(e)}")
        return {
            'statusCode': 500,
            'headers': {'Content-Type': 'application/json', 'Access-Control-Allow-Origin': '*'},
            'body': json.dumps({'error': f'Ошибка генерации: {str(e)}'}),
            'isBase64Encoded': False
        }

def generate_presentation(data: Dict[str, Any]) -> io.BytesIO:
    """Генерация презентации из шаблона"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data['title']
    subtitle.text = f"{data['university']}\n{data['faculty']}\n\nВыполнил: {data['author']}"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Актуальность темы"
    tf = body_shape.text_frame
    tf.text = data['relevance']
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Цель и задачи"
    tf = body_shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"Цель: {data['goal']}"
    p.level = 0
    
    for idx, task in enumerate(data['tasks'], 1):
        p = tf.add_paragraph()
        p.text = f"{idx}. {task}"
        p.level = 0
    
    for idx, task in enumerate(data['tasks'], 1):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"Задача {idx}"
        tf = body_shape.text_frame
        tf.text = task
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Выводы"
    tf = body_shape.text_frame
    
    for idx, conclusion in enumerate(data['conclusions'], 1):
        if idx == 1:
            tf.text = f"{idx}. {conclusion}"
        else:
            p = tf.add_paragraph()
            p.text = f"{idx}. {conclusion}"
            p.level = 0
    
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Спасибо за внимание!"
    subtitle.text = "Готов ответить на ваши вопросы"
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def generate_speech(data: Dict[str, Any]) -> io.BytesIO:
    """Генерация речи для защиты"""
    doc = Document()
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = DocPt(14)
    
    heading = doc.add_heading('РЕЧЬ ДЛЯ ЗАЩИТЫ', 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()
    
    p = doc.add_paragraph('Здравствуйте, уважаемые члены комиссии!')
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    doc.add_paragraph()
    
    doc.add_paragraph(f'Тема моей работы: «{data["title"]}».')
    
    doc.add_paragraph()
    
    doc.add_paragraph(f'Актуальность темы обусловлена тем, что {data["relevance"]}')
    
    doc.add_paragraph()
    
    doc.add_paragraph(f'Целью данной работы является {data["goal"]}')
    
    doc.add_paragraph()
    
    p = doc.add_paragraph('Для достижения цели были поставлены следующие задачи:')
    for idx, task in enumerate(data['tasks'], 1):
        doc.add_paragraph(f'{idx}. {task}', style='List Number')
    
    doc.add_paragraph()
    
    doc.add_paragraph('В ходе исследования были получены следующие результаты:')
    
    for idx, task in enumerate(data['tasks'], 1):
        doc.add_paragraph()
        doc.add_paragraph(f'В рамках {idx}-й задачи было выполнено: {task}.')
    
    doc.add_paragraph()
    
    doc.add_paragraph('На основании проведённого исследования можно сделать следующие выводы:')
    for idx, conclusion in enumerate(data['conclusions'], 1):
        doc.add_paragraph(f'{idx}. {conclusion}', style='List Number')
    
    doc.add_paragraph()
    
    doc.add_paragraph('Таким образом, цель работы достигнута, все задачи выполнены.')
    
    doc.add_paragraph()
    
    p = doc.add_paragraph('Спасибо за внимание! Готов ответить на ваши вопросы.')
    p.runs[0].bold = True
    
    doc.add_paragraph()
    doc.add_paragraph('_' * 60)
    doc.add_paragraph('Время выступления: ~3 минуты')
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_notes(data: Dict[str, Any]) -> str:
    """Генерация текстовой шпаргалки"""
    notes = f"""
═══════════════════════════════════════════════════════════
            ШПАРГАЛКА ДЛЯ ЗАЩИТЫ РАБОТЫ
═══════════════════════════════════════════════════════════

ТЕМА: {data['title']}

───────────────────────────────────────────────────────────
1. АКТУАЛЬНОСТЬ
───────────────────────────────────────────────────────────
{data['relevance']}

───────────────────────────────────────────────────────────
2. ЦЕЛЬ РАБОТЫ
───────────────────────────────────────────────────────────
{data['goal']}

───────────────────────────────────────────────────────────
3. ЗАДАЧИ
───────────────────────────────────────────────────────────
"""
    for idx, task in enumerate(data['tasks'], 1):
        notes += f"{idx}. {task}\n"
    
    notes += f"""
───────────────────────────────────────────────────────────
4. ВЫВОДЫ
───────────────────────────────────────────────────────────
"""
    for idx, conclusion in enumerate(data['conclusions'], 1):
        notes += f"{idx}. {conclusion}\n"
    
    notes += """
───────────────────────────────────────────────────────────
5. ТИПОВЫЕ ВОПРОСЫ И ОТВЕТЫ
───────────────────────────────────────────────────────────

❓ В чём актуальность вашей темы?
→ """ + data['relevance'] + """

❓ Какова цель вашей работы?
→ """ + data['goal'] + """

❓ Какие задачи вы решали?
→ """ + ", ".join(data['tasks']) + """

❓ Какие методы исследования использовали?
→ В работе применялись методы анализа, синтеза, 
  сравнительный метод, обобщение материала.

❓ Какова практическая значимость?
→ Результаты могут быть использованы в практической 
  деятельности предприятий и организаций данной сферы.

❓ Какие источники изучали?
→ Изучены работы отечественных и зарубежных авторов, 
  нормативные документы, статистические данные.

❓ Какие трудности возникли при написании?
→ Основная сложность заключалась в анализе большого 
  объёма материала и выборе наиболее релевантных данных.

❓ Почему выбрали именно эту тему?
→ Тема актуальна и соответствует моим профессиональным 
  интересам в данной области.

═══════════════════════════════════════════════════════════
           УДАЧИ НА ЗАЩИТЕ! 🎓
═══════════════════════════════════════════════════════════
"""
    return notes

def save_defense_kit_record(user_id: int, work_id: str):
    """Сохранение записи о создании пакета защиты"""
    database_url = os.environ.get('DATABASE_URL')
    conn = None
    try:
        conn = psycopg2.connect(database_url)
        cur = conn.cursor()
        
        cur.execute(
            """
            INSERT INTO t_p63326274_course_download_plat.defense_kits 
            (user_id, work_id, created_at) 
            VALUES (%s, %s, %s)
            """,
            (user_id, work_id, datetime.utcnow())
        )
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        if conn:
            conn.rollback()
            conn.close()
        raise